from __future__ import annotations

import io
import tempfile
from pathlib import Path

from app.logging_setup import logger


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------

def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract all text from a PowerPoint file."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))

        # Also grab speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Slide {i} — Speaker Notes]\n{notes}")

    return "\n\n".join(parts)


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract all text from a Word document."""
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    parts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    # Also grab tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n".join(parts)


def extract_text(filename: str, file_bytes: bytes) -> str:
    """Route to the correct extractor based on file extension."""
    ext = Path(filename).suffix.lower()
    if ext in (".pptx",):
        return extract_text_from_pptx(file_bytes)
    elif ext in (".docx",):
        return extract_text_from_docx(file_bytes)
    elif ext in (".txt", ".md"):
        return file_bytes.decode("utf-8", errors="replace")
    else:
        raise ValueError(f"Unsupported file type: {ext}. Use .pptx, .docx, or .txt")


# ---------------------------------------------------------------------------
# Synopsis generation
# ---------------------------------------------------------------------------

SYNOPSIS_PROMPT = """\
You are preparing an AI conference moderator for an upcoming talk. \
Below is the raw content extracted from the speaker's presentation materials.

Produce a BRIEFING DOCUMENT with these sections:

1. TOPIC: One sentence — what is this talk about?
2. KEY CLAIMS: 3-5 bullet points — what are the speaker's main arguments or findings?
3. TECHNICAL DETAILS: Any specific technologies, numbers, architectures, or results mentioned.
4. POTENTIAL WEAK SPOTS: 2-3 areas where the claims seem vague, unsubstantiated, \
   or where an audience of engineers would want more detail.
5. SUGGESTED ANGLES: 3 interesting directions a moderator could steer the Q&A \
   after the talk — think specific, provocative, audience-relevant.

Keep the briefing under 500 words. Be direct. This is an internal prep document, not a summary for the public.

EXTRACTED CONTENT:
{content}
"""


async def generate_synopsis(raw_text: str, llm_client, model: str) -> str:
    """Use the LLM to produce a structured briefing from raw document text."""
    # Truncate if the extracted text is very long
    truncated = raw_text[:8000]
    prompt = SYNOPSIS_PROMPT.format(content=truncated)

    try:
        resp = await llm_client.chat.completions.create(
            model=model,
            messages=[
                {"role": "user", "content": prompt},
            ],
            max_tokens=2048,
            temperature=0.4,  # more factual for synopsis
        )
        return resp.choices[0].message.content.strip()
    except Exception as exc:
        logger.error("Synopsis generation failed: %s", exc)
        # Fall back to a simple extractive summary
        lines = [l.strip() for l in raw_text.split("\n") if l.strip()]
        preview = "\n".join(lines[:30])
        return f"[Synopsis generation failed — showing raw extract]\n\n{preview}"
